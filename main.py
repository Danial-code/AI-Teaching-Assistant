import os
import re
import uuid
import json
import time
import fitz  # PyMuPDF
from pptx import Presentation
from langgraph.graph import StateGraph


# ---------------- LLM ---------------- #
try:
    from langchain_ollama import OllamaLLM
    llm = OllamaLLM(model="llama3")
    print("✓ Using langchain_ollama.OllamaLLM")
except ImportError:
    from langchain_community.llms import Ollama
    llm = Ollama(model="llama3")
    print("⚠️ Using deprecated langchain_community.Ollama")

# ---------------- Checkpointer ---------------- #
try:
    from langgraph.checkpoint.sqlite import SqliteSaver
    os.makedirs("state", exist_ok=True)
    memory = SqliteSaver.from_conn_string("state/checkpoints.db")
    print("✓ Using SqliteSaver (persistent)")
except Exception:
    from langgraph.checkpoint.memory import MemorySaver
    memory = MemorySaver()
    print("• SqliteSaver unavailable; using MemorySaver")

    STATE_DIR = "state"
    os.makedirs(STATE_DIR, exist_ok=True)

    def _state_path(tid): return os.path.join(STATE_DIR, f"{tid}.json")
    def _persist_to_disk(tid, state):
        try:
            with open(_state_path(tid), "w", encoding="utf-8") as f:
                json.dump(state, f, ensure_ascii=False, indent=2)
        except Exception:
            pass
    def _load_from_disk(tid):
        try:
            with open(_state_path(tid), "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            return None

# ---------------- File text extraction ---------------- #
def extract_text(file_path: str) -> str:
    try:
        if file_path.endswith(".pdf"):
            with fitz.open(file_path) as doc:
                return "\n".join(page.get_text() for page in doc)
        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = getattr(shape, "text", "").strip()
                        if t:
                            parts.append(t)
            return "\n".join(parts)
        else:
            return "⚠️ Unsupported file type."
    except Exception as e:
        return f"⚠️ Error reading file: {e}"

# ---------------- JSON helpers ---------------- #
def _extract_json_block(text: str) -> str | None:
    t = (text or "").strip()
    if t.startswith("```"):
        t = re.sub(r"^```(?:json)?\s*|\s*```$", "", t, flags=re.IGNORECASE | re.DOTALL).strip()
    try:
        json.loads(t); return t
    except Exception:
        pass
    stack = 0; start = None
    for i, ch in enumerate(t):
        if ch == "{":
            if stack == 0: start = i
            stack += 1
        elif ch == "}":
            stack -= 1
            if stack == 0 and start is not None:
                cand = t[start:i+1]
                try:
                    json.loads(cand); return cand
                except Exception:
                    continue
    return None

# ---------------- Grounding helpers ---------------- #
LANG_HINTS = {
    "java": ["public class_summary", "System.out.println", "static void", "static int", "String", "Scanner", "return", "void"],
    "python": ["def ", "print(", "import ", "lambda", "self:", "numpy", "pandas"],
    "javascript": ["function(", "=>", "console.log", "document.", "let ", "const ", "var "],
}
OFFTOPIC_FAMILIES = {
    "politics": ["election", "president", "policy", "parliament", "senate", "minister", "party"],
    "biology": ["cell", "organism", "DNA", "protein", "enzyme", "genome"],
    "history": ["war", "empire", "king", "dynasty", "revolution"],
    "geography": ["capital city", "continent", "latitude", "longitude"],
}
MATH_PATTERNS = [
    r"\b\d+\s*[\+\-\*/]\s*\d+",
    r"=\s*\d+",
    r"\bx\s*[\+\-\*/]",
    r"\bsolve for\b",
]

def detect_languages(lecture_text: str) -> set[str]:
    t = lecture_text.lower()
    langs = set()
    for lang, hints in LANG_HINTS.items():
        if any(h.lower() in t for h in hints):
            langs.add(lang)
    return langs

def lecture_has_math(lecture_text: str) -> bool:
    if re.search("|".join(MATH_PATTERNS), lecture_text, flags=re.IGNORECASE):
        return True
    return any(w in lecture_text.lower() for w in ["equation", "algebra", "derivative", "integral"])

def lecture_mentions_offtopic_family(lecture_text: str) -> set[str]:
    t = lecture_text.lower()
    hits = set()
    for fam, kws in OFFTOPIC_FAMILIES.items():
        if any(k.lower() in t for k in kws):
            hits.add(fam)
    return hits

def get_allowed_types(lecture_text: str) -> set[str]:
    allowed = {"mcq", "short_answer"}
    if detect_languages(lecture_text): allowed.add("coding")
    if lecture_has_math(lecture_text): allowed.add("math")
    return allowed

def content_in_scope(text: str, lecture_text: str, langs: set[str]) -> bool:
    t = (text or "").lower()
    lect = (lecture_text or "").lower()
    # block language drift
    for lang in ["python", "javascript", "js "]:
        if lang in t and lang not in langs: return False
    # block offtopic drift
    for fam, kws in OFFTOPIC_FAMILIES.items():
        if any(k in t for k in kws) and fam not in lecture_mentions_offtopic_family(lect): return False
    # require overlap with core lecture vocab
    core = ["method", "methods", "parameter", "argument", "return", "static", "void", "scope", "class_summary"]
    java_core = ["system.out.println", "public", "class_summary", "static", "void", "string", "int"]
    if any(w in lect for w in java_core) and any(w in t for w in java_core + core): return True
    return any(w in t and w in lect for w in core)

def enforce_grounding(quiz_json: dict, lecture_text: str) -> dict:
    allowed_types = get_allowed_types(lecture_text)
    langs = detect_languages(lecture_text)
    out = []
    for q in quiz_json.get("questions", []):
        qtype = (q.get("type") or "").lower().strip()
        if qtype not in allowed_types:
            continue
        blob = " ".join([q.get("question",""),
                         " ".join(q.get("options") or []),
                         q.get("answer","")])
        if content_in_scope(blob, lecture_text, langs):
            out.append(q)
    return {"questions": out}

# ---------------- Count & type parsing ---------------- #
def _normalize_qtype(s: str) -> str | None:
    s = (s or "").strip().lower()
    if s in {"mcq", "multiple choice", "multiple-choice"}: return "mcq"
    if s in {"short", "short answer", "short_answer"}: return "short_answer"
    if s in {"coding", "code"}: return "coding"
    if s in {"math", "equation"}: return "math"
    return None

def _parse_request_counts(instruction: str) -> tuple[int | None, dict]:
    instr = (instruction or "").lower()
    total = None
    m_total = re.search(r"(\d+)\s+questions?", instr)
    if m_total:
        total = int(m_total.group(1))
    per_type: dict[str, int] = {}
    for m in re.finditer(r"(\d+)\s*(mcq|multiple\s*choice|short\s*answer|coding|code|math|equations?)s?", instr):
        n = int(m.group(1))
        qt = _normalize_qtype(m.group(2))
        if qt:
            per_type[qt] = per_type.get(qt, 0) + n
    return total, per_type

# ---------------- LLM calls ---------------- #
def _llm_make_questions(lecture: str, request: str, allowed_types: set[str], n: int = 1, avoid_texts: list[str] | None = None) -> list[dict]:
    types_str = " | ".join(sorted(allowed_types)) if allowed_types else "mcq | short_answer"
    avoid_block = ""
    if avoid_texts:
        dedup = [s[:180] for s in avoid_texts if s]
        if dedup:
            avoid_block = "Avoid repeating or paraphrasing these existing questions:\n- " + "\n- ".join(dedup) + "\n"

    prompt = f"""
You will create {n} quiz question(s) STRICTLY from the LECTURE below.
Return ONLY JSON with a 'questions' array; do not include any prose.
{avoid_block}
LECTURE (Ground Truth):
{lecture}

REQUEST for new/replacement question(s):
{request}

ALLOWED TYPES: {types_str}

Return EXACTLY this schema:
{{
  "questions": [
    {{
      "type": "mcq" | "short_answer" | "coding" | "math",
      "question": "Question text?",
      "options": ["A) ...","B) ...","C) ...","D) ..."],
      "answer": "Answer text or correct option"
    }}
  ]
}}
"""
    raw = llm.invoke(prompt)
    blob = _extract_json_block(raw) or ""
    try:
        data = json.loads(blob) if blob else {}
    except Exception:
        data = {}
    qs = data.get("questions", [])
    return qs if isinstance(qs, list) else []

# ---------------- Exact collection (top-up until target) ---------------- #
def _collect_questions_exact(
    lecture: str,
    request_text: str,
    allowed_types: set[str],
    target_count: int,
    *,
    enforce_from_text: str,
    max_rounds: int = 12
) -> list[dict]:
    collected: list[dict] = []
    seen_pairs = set()
    rounds = 0

    def _append_unique(qs: list[dict]):
        for q in qs or []:
            key = (q.get("type",""), (q.get("question") or "").strip())
            if key[1] and key not in seen_pairs:
                seen_pairs.add(key)
                collected.append(q)

    while len(collected) < target_count and rounds < max_rounds:
        remaining = target_count - len(collected)
        avoid_texts = [q.get("question", "") for q in collected]
        req = f"{request_text}\n\nCreate {remaining} more question(s)."
        qs = _llm_make_questions(
            lecture,
            req,
            allowed_types=allowed_types,
            n=remaining,
            avoid_texts=avoid_texts
        )
        grounded_batch = enforce_grounding({"questions": qs}, enforce_from_text).get("questions", [])
        _append_unique(grounded_batch)
        rounds += 1

    return collected[:target_count]

# ---------------- Text ↔ JSON (preserve user edits) ---------------- #
def _parse_quiz_text_to_json(text: str) -> dict:

    if not text:
        return {"questions": []}

    lines = [ln.rstrip() for ln in text.splitlines()]
    questions = []
    cur = None

    def _flush():
        nonlocal cur, questions
        if cur:
            if cur["type"] == "mcq":
                cur["options"] = [o for o in (cur.get("options") or []) if o.strip()]
            else:
                cur["options"] = []
            questions.append(cur)
            cur = None

    q_re = re.compile(r"^\s*Question\s+(\d+)\s*\(([^)]+)\):\s*(.*)$", re.IGNORECASE)
    opt_re = re.compile(r"^\s*([A-Da-d])\)\s*(.*)$")
    ans_re = re.compile(r"^\s*Answer:\s*(.*)$", re.IGNORECASE)
    exp_re = re.compile(r"^\s*Expected\s*Answer:\s*(.*)$", re.IGNORECASE)

    for ln in lines:
        if not ln.strip():
            continue

        m_q = q_re.match(ln)
        if m_q:
            _flush()
            typ = m_q.group(2).strip().lower().replace(" ", "_")
            if typ not in {"mcq", "short_answer", "coding", "math"}:
                typ = "short_answer"
            cur = {"type": typ, "question": m_q.group(3).strip(), "options": [], "answer": ""}
            continue

        if cur is None:
            continue

        m_opt = opt_re.match(ln)
        if m_opt:
            if cur["type"] != "mcq":
                cur["type"] = "mcq"
                cur["options"] = []
            letter = m_opt.group(1).upper()
            rest = m_opt.group(2).strip()
            cur["options"].append(f"{letter}) {rest}")
            continue

        m_ans = ans_re.match(ln)
        if m_ans:
            cur["answer"] = m_ans.group(1).strip()
            continue

        m_exp = exp_re.match(ln)
        if m_exp:
            cur["answer"] = m_exp.group(1).strip()
            continue

        # continuation lines -> append to question
        if cur["question"]:
            cur["question"] += " " + ln.strip()
        else:
            cur["question"] = ln.strip()

    _flush()
    return {"questions": questions}

# ---------------- Follow-up merge (only ground NEW) ---------------- #
def apply_followup(prev: dict, followup: str, lecture_text: str) -> dict:

    allowed = get_allowed_types(lecture_text)
    out = {"questions": list(prev.get("questions", []))}
    follow = (followup or "").strip().lower()
    if not follow:
        return out

    # ------- change question N [to <type>] -------
    for m in re.finditer(r"change\s+question\s+(\d+)(?:\s+to\s+(mcq|short\s*answer|coding|code|math))?", follow):
        idx = int(m.group(1)) - 1
        target_type = (m.group(2) or "").replace(" ", "_")
        if target_type == "code":
            target_type = "coding"
        if 0 <= idx < len(out["questions"]):
            req_type = target_type or out["questions"][idx].get("type", "mcq")
            if req_type in allowed:
                avoid = [q.get("question","") for q in out["questions"]]
                req = f"Create 1 '{req_type}' question from the lecture covering a different angle than existing Q{idx+1}."
                new_qs = _llm_make_questions(lecture_text, req, allowed_types={req_type}, n=1, avoid_texts=avoid)
                if new_qs:
                    grounded = enforce_grounding({"questions": new_qs}, lecture_text).get("questions", [])
                    if grounded:
                        out["questions"][idx] = grounded[0]


    remaining = follow

    # ------- add N <type> [question|questions] (optional 'more') -------
    typed_pat = re.compile(
        r"add\s+(\d+)\s+(?:more\s+)?(mcq|short\s*answer|coding|code|math)(?:\s+questions?)?",
        re.IGNORECASE
    )
    for m in typed_pat.finditer(remaining):
        n = int(m.group(1))
        t = m.group(2).replace(" ", "_").lower()
        if t == "code":
            t = "coding"
        if t in allowed and n > 0:
            batch = _collect_questions_exact(
                lecture_text,
                f"Strictly from the lecture, generate '{t}' questions only.",
                allowed_types={t},
                target_count=n,
                enforce_from_text=lecture_text,
            )
            out["questions"].extend(batch)
        # mask handled phrase to avoid double handling by generic rule
        span = m.span()
        remaining = remaining[:span[0]] + (" " * (span[1] - span[0])) + remaining[span[1]:]

    # ------- add N (more) question(s) (any allowed type) -------
    any_pat = re.compile(r"add\s+(\d+)\s+(?:more\s+)?questions?", re.IGNORECASE)
    for m in any_pat.finditer(remaining):
        n = int(m.group(1))
        if n > 0:
            batch = _collect_questions_exact(
                lecture_text,
                "Strictly from the lecture, generate questions of any allowed type.",
                allowed_types=allowed,
                target_count=n,
                enforce_from_text=lecture_text,
            )
            out["questions"].extend(batch)

    return out

# ---------------- LangGraph nodes ---------------- #
def summarize_lecture(state: dict) -> dict:
    lecture = state["lecture_text"]
    summary = llm.invoke(f"Summarize for quiz generation:\n\n{lecture}")
    return {**state, "summary": summary}

def generate_quiz_json(state: dict) -> dict:
    lecture_text = state.get("lecture_text","")
    instruction = state["instruction"]
    allowed = get_allowed_types(lecture_text)

    total_req, per_type_req = _parse_request_counts(instruction)
    questions: list[dict] = []

    if per_type_req:
        for qtype, count in per_type_req.items():
            if qtype in allowed and count > 0:
                batch = _collect_questions_exact(
                    lecture_text,
                    f"Strictly from the lecture, generate '{qtype}' questions only.",
                    allowed_types={qtype},
                    target_count=count,
                    enforce_from_text=lecture_text,
                )
                questions.extend(batch)

        if total_req and len(questions) < total_req:
            topup = _collect_questions_exact(
                lecture_text,
                "Strictly from the lecture, generate additional questions of any allowed type.",
                allowed_types=allowed,
                target_count=total_req - len(questions),
                enforce_from_text=lecture_text,
            )
            questions.extend(topup)
    else:
        target = total_req if (total_req and total_req > 0) else 5
        questions = _collect_questions_exact(
            lecture_text,
            "Strictly from the lecture, generate questions of any allowed type.",
            allowed_types=allowed,
            target_count=target,
            enforce_from_text=lecture_text,
        )

    strict = enforce_grounding({"questions": questions}, lecture_text)
    return {**state, "quiz_json": strict}

# ---------------- Graph ---------------- #
workflow = StateGraph(dict)
workflow.add_node("summarize", summarize_lecture)
workflow.add_node("quiz", generate_quiz_json)
workflow.add_edge("summarize", "quiz")
workflow.set_entry_point("summarize")
workflow.set_finish_point("quiz")
graph = workflow.compile(checkpointer=memory)

# ---------------- Formatter ---------------- #
def format_quiz_text(data: dict) -> str:
    lines = []
    for idx, q in enumerate(data.get("questions", []), 1):
        typ = (q.get("type","mcq") or "mcq").lower()
        typ_title = typ.title().replace("_", " ")
        lines.append(f"Question {idx} ({typ_title}): {q.get('question','')}")
        if typ == "mcq":
            for opt in q.get("options", []): lines.append(opt)
            if q.get("answer"): lines.append(f"Answer: {q['answer']}")
        else:
            if q.get("answer"): lines.append(f"Expected Answer: {q['answer']}")
        lines.append("")
    return "\n".join(lines).strip()

# ---------------- Session helpers ---------------- #
def _cfg(tid): return {"configurable": {"thread_id": tid, "checkpoint_ns": "ai_ta", "checkpoint_id": tid}}

def _save_session(tid, *, lecture_text, instruction, quiz_json, file_path):
    config = _cfg(tid); now = int(time.time())
    try: ckpt = memory.get(config) or {}
    except Exception: ckpt = {}
    state = ckpt.get("state") or {}
    state.update({
        "lecture_text": lecture_text,
        "instruction": instruction,
        "quiz_json": quiz_json,
        "file_path": file_path,
        "updated_at": now,
    })
    checkpoint = {"id": tid, "v": ckpt.get("v") or 1, "ts": now, "state": state, "channel_values": {}}
    try: memory.put(config, checkpoint, {}, {})
    except Exception: pass
    if "SqliteSaver" not in memory.__class__.__name__:
        try: _persist_to_disk(tid, state)
        except Exception: pass

def _load_from_disk(tid):
    try:
        with open(os.path.join("state", f"{tid}.json"), "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return None

def _load_session(tid):
    config = _cfg(tid)
    try:
        ckpt = memory.get(config)
        if ckpt and "state" in ckpt:
            return ckpt["state"]
    except Exception:
        pass
    if "SqliteSaver" not in memory.__class__.__name__:
        return _load_from_disk(tid)
    return None

# ---------------- Public API ---------------- #
def generate_quiz_stream(lecture_text: str, instruction: str, file_path: str):
    tid = str(uuid.uuid4())
    result = graph.invoke({"lecture_text": lecture_text, "instruction": instruction}, config=_cfg(tid))
    quiz_json = result.get("quiz_json", {"questions": []})
    _save_session(tid, lecture_text=lecture_text, instruction=instruction, quiz_json=quiz_json, file_path=file_path)
    text = format_quiz_text(quiz_json) or "No questions generated."
    return text, tid

def regenerate_quiz_stream(edited_quiz: str, followup: str, tid: str) -> str:
    if not tid:
        return "No active quiz context found for this session.\nTip: Generate a new quiz first."
    saved = _load_session(tid)
    if not saved:
        return "No active quiz context found for this session.\nTip: Generate a new quiz first."

    prev_json = saved.get("quiz_json") or {"questions": []}

    # 1) Merge user's manual edits (from typingBox) by index
    parsed = _parse_quiz_text_to_json(edited_quiz)
    merged = {"questions": list(prev_json.get("questions", []))}

    if parsed.get("questions"):
        new_qs = parsed["questions"]
        for i in range(min(len(merged["questions"]), len(new_qs))):
            merged["questions"][i] = {
                "type": (new_qs[i].get("type") or merged["questions"][i].get("type") or "short_answer"),
                "question": new_qs[i].get("question", "").strip() or merged["questions"][i].get("question",""),
                "options": list(new_qs[i].get("options") or (merged["questions"][i].get("options") or [])),
                "answer": new_qs[i].get("answer", "").strip() or merged["questions"][i].get("answer",""),
            }
        if len(new_qs) > len(merged["questions"]):
            merged["questions"].extend(new_qs[len(merged["questions"]):])
    else:
        merged = prev_json

    # 2) Apply follow-up (only new items are grounded)
    new_json = apply_followup(merged, followup, saved.get("lecture_text",""))

    # 3) Save & return
    _save_session(
        tid,
        lecture_text=saved.get("lecture_text",""),
        instruction=saved.get("instruction",""),
        quiz_json=new_json,
        file_path=saved.get("file_path",""),
    )
    return format_quiz_text(new_json) or format_quiz_text(merged)
