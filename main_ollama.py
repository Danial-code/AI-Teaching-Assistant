from __future__ import annotations

import os
import time
import re
import uuid
from typing import Optional, Tuple

import ollama

# Optional document parsers
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from langchain_ollama import OllamaLLM
    _LLM_SUMMARY = OllamaLLM(model=os.getenv("OLLAMA_MODEL", "llama3"), temperature=0.10)
except Exception:
    _LLM_SUMMARY = None



# CONFIG


OLLAMA_MODEL = os.getenv("OLLAMA_MODEL", "llama3")
MAX_SOURCE_CHARS = 22000

# In-memory context for strict grounding during regeneration.
# Maps thread_id -> {"lecture_text": str, "quiz_text": str}
_SESSION_STORE: dict[str, dict[str, str]] = {}
_SESSION_STORE_MAX = 30

def _session_put(tid: str, lecture_text: str, quiz_text: str) -> None:
    if not tid:
        return
    # simple FIFO eviction
    if tid not in _SESSION_STORE and len(_SESSION_STORE) >= _SESSION_STORE_MAX:
        oldest = next(iter(_SESSION_STORE.keys()))
        _SESSION_STORE.pop(oldest, None)
    _SESSION_STORE[tid] = {"lecture_text": lecture_text or "", "quiz_text": quiz_text or ""}


def _session_get_lecture(tid: str) -> str:
    return (_SESSION_STORE.get(tid) or {}).get("lecture_text", "")


def _clip(text: str, limit: int = MAX_SOURCE_CHARS) -> str:
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    head = text[: int(limit * 0.75)]
    tail = text[-int(limit * 0.25):]
    return head + "\n\n[...TRUNCATED...]\n\n" + tail


def _chat(messages, temperature=0.2, max_tokens=1200, retries=2) -> str:
    last_err = None
    for i in range(retries + 1):
        try:
            res = ollama.chat(
                model=OLLAMA_MODEL,
                messages=messages,
                options={
                    "temperature": temperature,
                    "num_predict": max_tokens,
                },
                stream=False,
            )
            return res.get("message", {}).get("content", "").strip()
        except Exception as e:
            last_err = e
            time.sleep(0.6 * (i + 1))
    return f"⚠️ Ollama error: {type(last_err).__name__}: {last_err}"


def extract_text(path: str) -> str:
    if not path:
        return ""

    ext = os.path.splitext(path)[1].lower()

    if ext == ".pdf":
        if not fitz:
            return "⚠️ PyMuPDF not installed"
        try:
            with fitz.open(path) as doc:
                return "\n".join(page.get_text() for page in doc)
        except Exception as e:
            return f"⚠️ PDF read error: {e}"

    if ext == ".pptx":
        if not Presentation:
            return "⚠️ python-pptx not installed"
        try:
            prs = Presentation(path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = shape.text.strip()
                        if t:
                            parts.append(t)
            return "\n".join(parts)
        except Exception as e:
            return f"⚠️ PPTX read error: {e}"

    return "⚠️ Unsupported file type"

# QUIZ GENERATION

SYS_QUIZ = (
    "You are an academic quiz generator for a university. "
    "You MUST generate questions ONLY from the provided SOURCE MATERIAL. "
    "Do NOT infer, generalize, or add concepts, equations, or examples not explicitly present. "
    "If something is not stated in the material, you must not ask about it."
)

FORMAT_RULES = """
FORMAT RULES (STRICT — MUST FOLLOW):
- Output MUST be a single numbered list from Question 1 to Question N
- Each question must start EXACTLY with: Question N (Type):
- Type must be one of: MCQ, Short_Answer, True_False, Coding, Math
- DO NOT group questions under headings
- MCQ must have EXACTLY 4 options: A) B) C) D)
- EVERY question MUST include an Answer line
- Answers MUST be concrete (no 'not specified', no 'depends', no 'not stated')
- DO NOT include notes, explanations, comments, warnings, or meta text
- DO NOT mention missing information
- IF there are limited concepts, you MUST still generate the requested number of questions by rephrasing, varying question type, or changing perspective, but ONLY using explicitly stated concepts from the source material.
"""

DEFAULT_QUIZ_INSTRUCTION = (
    "Create a 10-question quiz from the lecture. "
    "Mix MCQ and Short_Answer. "
    "Include Coding or Math only if the lecture contains them."
)


def generate_quiz_stream(
    lecture_text: str,
    instruction: str,
    source_path: Optional[str]
) -> Tuple[str, str]:

    thread_id = uuid.uuid4().hex
    lecture_text = _clip(lecture_text)
    instruction = instruction.strip() if instruction else DEFAULT_QUIZ_INSTRUCTION

    m_n = re.search(r"\b(\d+)\s+questions?\b", instruction, flags=re.IGNORECASE)
    target_n = int(m_n.group(1)) if m_n else 10
    target_n = max(1, min(30, target_n))

    prompt = f"""
SOURCE MATERIAL (only allowed knowledge):
{lecture_text}

INSTRUCTION:
{instruction}

HARD REQUIREMENT:
- Output EXACTLY {target_n} questions.

{FORMAT_RULES}

STRICT CONSTRAINT:
- Every question and answer must be directly supported by the SOURCE MATERIAL above.
- If the material does not mention a concept (e.g., multimodality, extra equations), you must NOT include it.
""".strip()

    messages = [
        {"role": "system", "content": SYS_QUIZ},
        {"role": "user", "content": prompt},
    ]

    quiz_text = _chat(messages, temperature=0.15, max_tokens=1500)

    def _count_questions(txt: str) -> int:
        t = txt or ""
        # Accept both 'Question 1 (..):' and '1. Question 1 (..):'
        return len(re.findall(r"^\s*(?:\d+[\.)]\s*)?Question\s+\d+\s*\(", t, flags=re.IGNORECASE | re.MULTILINE))

    cur_n = _count_questions(quiz_text)
    if cur_n < target_n:
        base_quiz = (quiz_text or "").strip()
        rounds = 0
        while cur_n < target_n and rounds < 4:
            missing = target_n - cur_n
            start_n = cur_n + 1
            end_n = cur_n + missing

            follow_prompt = f"""
SOURCE MATERIAL (only allowed knowledge):
{lecture_text}

TASK:
Append {missing} new question(s) to the end of the quiz below.

RULES:
- Keep ALL existing questions unchanged.
- Append ONLY the new questions at the end.
- The new questions MUST be numbered from Question {start_n} to Question {end_n}.
- Follow the same strict format and include Answer for each.
- Do NOT add headings, notes, warnings, or extra text.

CURRENT QUIZ:
{base_quiz}
""".strip()

            extra = _chat(
                [
                    {"role": "system", "content": SYS_QUIZ},
                    {"role": "user", "content": follow_prompt},
                ],
                temperature=0.10,
                max_tokens=1600,
            )

            base_quiz = (extra or base_quiz).strip()
            cur_n = _count_questions(base_quiz)
            rounds += 1

        quiz_text = base_quiz

    _session_put(thread_id, lecture_text, quiz_text)
    return quiz_text, thread_id


def regenerate_quiz_stream(edited_quiz: str, followup: str, thread_id: str) -> str:

    followup = (followup or "").strip()

    if not followup:
        return edited_quiz

    lecture_text = _session_get_lecture(thread_id)
    lecture_text = _clip(lecture_text) if lecture_text else ""

    m = re.search(r"\b(?:change|replace)\s+question\s+(\d+)\b", followup, flags=re.IGNORECASE)
    replace_qnum = int(m.group(1)) if m else None

    SYSTEM_RULES = """
You are editing an EXISTING quiz.

STRICT RULES (MANDATORY):
- You MUST keep ALL unchanged questions EXACTLY as they are (word-for-word).
- You MUST NOT introduce any outside knowledge.
- Any new or replaced question MUST be directly supported by the SOURCE MATERIAL.
- Do NOT add new terminology, equations, or examples not present in the SOURCE MATERIAL.
- Keep the same question numbering unless adding/removing.
- If asked to ADD a question → append it at the end.
- If asked to DELETE a question → remove ONLY that question.

REPLACE / CHANGE RULE (IMPORTANT):
- If the user requests to CHANGE/REPLACE question N, replace the ENTIRE question N block (question text + options + Answer line).
- The replacement must be substantially different, but still strictly from the SOURCE MATERIAL.

OUTPUT:
- Output the FULL updated quiz (not a diff)
""".strip()

    prompt = (
        f"SOURCE MATERIAL (only allowed knowledge):\n{lecture_text}\n\n" if lecture_text else ""
    ) + (
        f"{SYSTEM_RULES}\n\n"
        f"USER INSTRUCTION:\n{followup}\n"
        + (f"\nSPECIAL INSTRUCTION: Replace ONLY Question {replace_qnum} entirely (rewrite it from scratch, including options and Answer).\n" if replace_qnum else "")
        + "\nCURRENT QUIZ:\n----------------\n"
        + f"{edited_quiz}\n"
        + "----------------\n\n"
        + "Return the updated quiz that reflects ONLY the requested change."
    )

    messages = [
        {"role": "system", "content": "You are a careful academic quiz editor."},
        {"role": "user", "content": prompt},
    ]

    updated = _chat(messages, temperature=0.08, max_tokens=1600)
    _session_put(thread_id, lecture_text or "", updated)
    return updated

# SUMMARIZATION


def summarize_lecture(transcript: str, lecture_text: str = "") -> str:
    transcript = _clip(transcript, 16000)
    lecture_text = _clip(lecture_text, 16000)

    prompt = f"""
You are an academic lecture note-taker.

ABSOLUTE RULES (MANDATORY):
- You MUST ONLY use information that explicitly appears in the LECTURE NOTES or TRANSCRIPT below.
- You MUST NOT invent topics, concepts, formulas, titles, or examples.
- If a section is not clearly supported by the inputs, write: "Not covered in the lecture."
- If the lecture topic cannot be confidently inferred, set the title to: "Lecture Title: (Not stated)".
- NEVER summarize generic subjects (e.g., calculus, physics, statistics) unless they clearly appear in the inputs.
- NEVER add exam advice that is not explicitly stated.

LECTURE NOTES (ground truth):
{lecture_text}

TRANSCRIPT (ground truth):
{transcript}

Write the notes in EXACTLY this structure:

# Lecture Title

## Short Overview

## Key Concepts
- bullet points ONLY from the inputs

## Important Details / Examples
- bullet points ONLY from the inputs

## Formulas / Definitions (if any)
- bullet points ONLY from the inputs, otherwise write: "Not covered in the lecture."

## Summary Tips for Exam
- bullet points ONLY from the inputs, otherwise write: "Not covered in the lecture."

OUTPUT ONLY the notes. Do NOT add explanations or comments.
""".strip()

    try:
        if _LLM_SUMMARY:
            return _LLM_SUMMARY.invoke(prompt).strip()
        else:
            return _chat(
                [
                    {"role": "system", "content": "You summarize university lectures."},
                    {"role": "user", "content": prompt},
                ],
                temperature=0.10,
                max_tokens=1200,
            )
    except Exception as e:
        return f"⚠️ Summary error: {e}"